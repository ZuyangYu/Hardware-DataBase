"""Result snapshot/export contracts.

These tests intentionally exercise the durable store and real renderers.  A
future refactor may change the backing queue or XML layout, but it must not
change the user-visible guarantees covered here.
"""

from __future__ import annotations

import io
import zipfile
from datetime import datetime, timedelta, timezone

import pytest

import src.settings
from src.result_exports.models import ResultEnvelope, normalize_export_format
from src.result_exports.renderers import render_result
from src.result_exports.store import ResultExportStore
from src.result_exports.worker import ResultExportWorker


def _envelope() -> ResultEnvelope:
    return ResultEnvelope(
        title="器件检索结果",
        query="查询电源芯片",
        answer="知识库中找到 1 条匹配记录。",
        tables=[
            {
                "name": "检索结果",
                "columns": ["型号", "数量"],
                "rows": [["TPS62130", 2]],
            }
        ],
        citations=[
            {
                "index": 1,
                "title": "power.xlsx",
                "locator": "Sheet1!A2",
                "excerpt": "TPS62130",
            }
        ],
        metadata={"knowledge_base": "hardware"},
    )


def test_snapshot_and_export_job_are_idempotent_and_can_publish_artifact(tmp_path):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    first_snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    replay_snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )

    assert replay_snapshot.snapshot_id == first_snapshot.snapshot_id
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=first_snapshot.snapshot_id,
        format="md",
        content_shape="report",
        client_request_id="export-1",
    )
    assert store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=first_snapshot.snapshot_id,
        format="md",
        content_shape="report",
        client_request_id="export-1",
    ).export_job_id == job.export_job_id

    claimed = store.claim(job.export_job_id, "export-worker-a", lease_seconds=30)
    assert claimed is not None
    artifact = store.publish_artifact(
        job_id=claimed.export_job_id,
        worker_id="export-worker-a",
        lease_token=claimed.lease_token,
        content=b"# result\n",
        filename="result.md",
        mime_type="text/markdown; charset=utf-8",
        preview={"format": "md"},
    )
    completed = store.get_export_job(7, job.export_job_id)

    assert completed is not None and completed.status == "succeeded"
    assert completed.artifact_id == artifact.artifact_id
    assert store.read_artifact(7, artifact.artifact_id) == b"# result\n"


def test_renderers_keep_answer_tables_and_citations_in_markdown_and_xlsx():
    envelope = _envelope()

    markdown = render_result(envelope, "md")
    assert markdown.mime_type == "text/markdown; charset=utf-8"
    assert "知识库中找到 1 条匹配记录" in markdown.content.decode("utf-8")
    assert "| 型号 | 数量 |" in markdown.content.decode("utf-8")
    assert "power.xlsx" in markdown.content.decode("utf-8")

    workbook = render_result(envelope, "xlsx")
    assert workbook.mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    with zipfile.ZipFile(io.BytesIO(workbook.content)) as archive:
        names = set(archive.namelist())
        assert "xl/workbook.xml" in names
        xml = "\n".join(archive.read(name).decode("utf-8") for name in names if name.startswith("xl/worksheets/"))
    assert "TPS62130" in xml
    assert "知识库中找到 1 条匹配记录" in xml


def test_renderers_support_word_pdf_and_powerpoint_with_citations_and_charts():
    envelope = _envelope()

    word = render_result(envelope, "docx")
    assert word.mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    assert word.extension == "docx"
    assert word.content.startswith(b"PK")
    from docx import Document

    document = Document(io.BytesIO(word.content))
    assert any("知识库中找到 1 条匹配记录" in paragraph.text for paragraph in document.paragraphs)
    assert word.preview["citation_count"] == 1

    pdf = render_result(envelope, "pdf")
    assert pdf.mime_type == "application/pdf"
    assert pdf.extension == "pdf"
    assert pdf.content.startswith(b"%PDF-")
    assert pdf.content.rstrip().endswith(b"%%EOF")
    assert pdf.preview["citation_count"] == 1

    power_point = render_result(
        envelope,
        "pptx",
        render_options={"theme": "dark", "include_charts": True},
    )
    assert power_point.mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    assert power_point.extension == "pptx"
    assert power_point.content.startswith(b"PK")
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(power_point.content))
    assert len(presentation.slides) >= 3
    assert any(shape.has_chart for slide in presentation.slides for shape in slide.shapes)
    assert power_point.preview["theme"] == "dark"


def test_export_format_aliases_normalize_to_stable_artifact_extensions():
    assert normalize_export_format("markdown") == "md"
    assert normalize_export_format("excel") == "xlsx"
    assert normalize_export_format("word") == "docx"
    assert normalize_export_format("woed") == "docx"
    assert normalize_export_format(".pdf") == "pdf"
    assert normalize_export_format("powerpoint") == "pptx"


def test_unknown_export_format_fails_before_a_job_is_created():
    with pytest.raises(ValueError, match="unsupported export format"):
        render_result(_envelope(), "txt")


def test_export_worker_renders_queued_job_and_records_downloadable_artifact(tmp_path):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="xlsx",
        content_shape="data",
        client_request_id="export-xlsx-1",
    )

    worker = ResultExportWorker(store=store, worker_id="export-worker-a")
    assert worker.run_once() is True

    completed = store.get_export_job(7, job.export_job_id)
    assert completed is not None and completed.status == "succeeded"
    assert completed.artifact_id
    artifact = store.get_artifact(7, completed.artifact_id)
    assert artifact is not None
    assert artifact.filename.endswith(".xlsx")
    assert store.read_artifact(7, artifact.artifact_id).startswith(b"PK")


@pytest.mark.parametrize("format", ["md", "xlsx", "docx", "pdf", "pptx"])
def test_export_worker_generates_a_valid_artifact_for_every_released_format(tmp_path, format):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id=f"turn-{format}",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format=format,
        client_request_id=f"all-formats-{format}",
    )

    assert ResultExportWorker(store=store, worker_id=f"worker-{format}").run_once() is True
    completed = store.get_export_job(7, job.export_job_id)
    assert completed is not None and completed.status == "succeeded" and completed.artifact_id
    artifact = store.get_artifact(7, completed.artifact_id)
    assert artifact is not None and artifact.filename.endswith(f".{format}")
    assert store.read_artifact(7, artifact.artifact_id)


def test_manual_retry_resets_the_attempt_budget_for_a_dead_letter_job(tmp_path):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="md",
        client_request_id="retry-dead-letter",
        max_attempts=1,
    )
    claimed = store.claim(job.export_job_id, "export-worker-a", lease_seconds=30)
    assert claimed is not None
    dead_letter = store.fail(
        claimed.export_job_id,
        "export-worker-a",
        claimed.lease_token,
        "transient renderer error",
    )
    assert dead_letter.status == "dead_letter"

    retried = store.retry(7, job.export_job_id)
    assert retried is not None
    assert retried.status == "queued"
    assert retried.attempt == 0
    assert ResultExportWorker(store=store, worker_id="export-worker-b").run_once() is True
    completed = store.get_export_job(7, job.export_job_id)
    assert completed is not None and completed.status == "succeeded"


def test_artifact_extension_must_match_the_claimed_job_format(tmp_path):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="md",
        client_request_id="format-mismatch",
    )
    claimed = store.claim(job.export_job_id, "export-worker-a", lease_seconds=30)
    assert claimed is not None
    with pytest.raises(ValueError, match="does not match export job"):
        store.publish_artifact(
            job_id=claimed.export_job_id,
            worker_id="export-worker-a",
            lease_token=claimed.lease_token,
            content=b"not xlsx",
            filename="result.xlsx",
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )


def test_artifact_signature_and_mime_must_match_the_claimed_job_format(tmp_path):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="xlsx",
        client_request_id="signature-mismatch",
    )
    claimed = store.claim(job.export_job_id, "export-worker-a", lease_seconds=30)
    assert claimed is not None
    with pytest.raises(ValueError, match="signature"):
        store.publish_artifact(
            job_id=claimed.export_job_id,
            worker_id="export-worker-a",
            lease_token=claimed.lease_token,
            content=b"not an office package",
            filename="result.xlsx",
            mime_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        )


def test_render_options_are_job_scoped_while_the_snapshot_stays_immutable(tmp_path):
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="md",
        client_request_id="render-options",
        options={"render_title": "定制报告", "include_citations": False},
    )

    assert ResultExportWorker(store=store, worker_id="export-worker-a").run_once() is True
    artifact = store.get_artifact(7, store.get_export_job(7, job.export_job_id).artifact_id)
    assert artifact is not None
    content = store.read_artifact(7, artifact.artifact_id).decode("utf-8")
    assert content.startswith("# 定制报告")
    assert "参考来源" not in content
    assert store.get_snapshot(7, snapshot.snapshot_id).envelope.title == "器件检索结果"


def test_claim_respects_global_running_job_quota(tmp_path, monkeypatch):
    monkeypatch.setattr(src.settings, "RESULT_EXPORT_MAX_RUNNING_JOBS", 1, raising=False)
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    first = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="md",
        client_request_id="quota-1",
    )
    second = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="pdf",
        client_request_id="quota-2",
    )

    assert store.claim(first.export_job_id, "worker-a", lease_seconds=30) is not None
    assert store.claim(second.export_job_id, "worker-b", lease_seconds=30) is None
    store.cancel(7, first.export_job_id)
    assert store.claim(second.export_job_id, "worker-b", lease_seconds=30) is not None


def test_expired_artifacts_are_removed_without_removing_export_history(tmp_path, monkeypatch):
    monkeypatch.setattr(src.settings, "RESULT_EXPORT_RETENTION_DAYS", 1, raising=False)
    store = ResultExportStore(str(tmp_path / "auth.db"), storage_dir=str(tmp_path / "exports"))
    snapshot = store.create_snapshot(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        turn_id="turn-1",
        envelope=_envelope(),
    )
    job = store.create_export_job(
        owner_user_id=7,
        tenant_id="default",
        session_id=3,
        snapshot_id=snapshot.snapshot_id,
        format="md",
        client_request_id="retention-1",
    )
    assert ResultExportWorker(store=store, worker_id="worker-a").run_once() is True
    completed = store.get_export_job(7, job.export_job_id)
    assert completed is not None and completed.artifact_id
    artifact = store.get_artifact(7, completed.artifact_id)
    assert artifact is not None and artifact.expires_at
    path = store.storage_dir / artifact.storage_ref
    assert path.exists()

    removed = store.cleanup_expired(now=datetime.now(timezone.utc) + timedelta(days=2))

    assert [item.artifact_id for item in removed] == [artifact.artifact_id]
    assert not path.exists()
    assert store.get_artifact(7, artifact.artifact_id) is None
    retained_job = store.get_export_job(7, job.export_job_id)
    assert retained_job is not None and retained_job.status == "succeeded"
