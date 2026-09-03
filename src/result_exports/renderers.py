"""Deterministic renderers for result exports.

The worker receives an immutable :class:`ResultEnvelope` and a small set of
presentation options.  The renderers do not execute user supplied code or
follow external links; every output is built from that envelope only.
"""

from __future__ import annotations

import io
import re
import zipfile
from dataclasses import replace
from datetime import date, datetime
import math
from typing import Any
from xml.sax.saxutils import escape

from src.result_exports.models import RenderedResult, ResultEnvelope, normalize_content_shape, normalize_export_format


MARKDOWN_MIME = "text/markdown; charset=utf-8"
XLSX_MIME = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
PDF_MIME = "application/pdf"
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
_MAX_ROWS = 10_000
_MAX_COLUMNS = 100
_MAX_CITATIONS = 200


def _text(value: Any) -> str:
    if value is None:
        return ""
    if isinstance(value, bool):
        return "是" if value else "否"
    return str(value)


def _markdown_cell(value: Any) -> str:
    return _text(value).replace("|", "\\|").replace("\r", " ").replace("\n", "<br>")


def _render_markdown(envelope: ResultEnvelope, content_shape: str) -> RenderedResult:
    if content_shape == "raw":
        raw = envelope.answer.rstrip() or envelope.query.rstrip()
        return RenderedResult(
            content=(raw + "\n").encode("utf-8"),
            mime_type=MARKDOWN_MIME,
            extension="md",
            preview={"format": "md", "text_preview": raw[:4000]},
        )
    lines = [f"# {envelope.title}", ""]
    if content_shape == "report" and envelope.query:
        lines.extend([f"**查询：** {envelope.query}", ""])
    if content_shape == "report" and envelope.answer:
        lines.extend([envelope.answer.rstrip(), ""])
    if content_shape == "report" and envelope.footer:
        lines.extend(["## 检索概览", "", envelope.footer.rstrip(), ""])
    for table in envelope.tables if content_shape != "raw" else []:
        name = str(table.get("name") or "检索结果")
        columns = [str(column) for column in (table.get("columns") or [])]
        rows = table.get("rows") or []
        if not columns:
            continue
        lines.extend([f"## {name}", "", "| " + " | ".join(_markdown_cell(column) for column in columns) + " |"])
        lines.append("| " + " | ".join("---" for _ in columns) + " |")
        for row in rows[:_MAX_ROWS]:
            values = row if isinstance(row, (list, tuple)) else [row]
            lines.append("| " + " | ".join(_markdown_cell(values[index] if index < len(values) else "") for index in range(len(columns))) + " |")
        lines.append("")
    if envelope.citations and content_shape != "raw":
        lines.extend(["## 参考来源", ""])
        for citation in envelope.citations[:_MAX_CITATIONS]:
            title = _text(citation.get("title") or "参考来源")
            locator = _text(citation.get("locator"))
            evidence_id = _text(citation.get("evidence_id"))
            details = " · ".join(item for item in (locator, f"证据 {evidence_id}" if evidence_id else "") if item)
            suffix = f"（{details}）" if details else ""
            lines.append(f"{citation.get('index', '')}. **{title}**{suffix}")
            excerpt = _text(citation.get("excerpt"))
            if excerpt:
                lines.append(f"   > {excerpt.replace(chr(10), ' ')}")
        lines.append("")
    return RenderedResult(
        content=("\n".join(lines).rstrip() + "\n").encode("utf-8"),
        mime_type=MARKDOWN_MIME,
        extension="md",
        preview={"format": "md", "text_preview": "\n".join(lines)[:4000]},
    )


def _column_name(index: int) -> str:
    result = ""
    value = index + 1
    while value:
        value, remainder = divmod(value - 1, 26)
        result = chr(65 + remainder) + result
    return result


def _safe_sheet_name(value: str, used: set[str]) -> str:
    base = re.sub(r"[\\/*?:\[\]]", "-", value or "结果").strip() or "结果"
    base = base[:31]
    name = base
    suffix = 1
    while name in used:
        suffix_text = f"-{suffix}"
        name = f"{base[:31 - len(suffix_text)]}{suffix_text}"
        suffix += 1
    used.add(name)
    return name


def _excel_date_serial(value: date | datetime) -> float:
    """Convert a Python date to the default 1900-based Excel serial."""

    if isinstance(value, datetime):
        value = value.replace(tzinfo=None)
        day = value.date()
        fraction = (
            value.hour * 3600 + value.minute * 60 + value.second + value.microsecond / 1_000_000
        ) / 86_400
    else:
        day = value
        fraction = 0.0
    return float((day - date(1899, 12, 30)).days) + fraction


def _cell_xml(reference: str, value: Any) -> str:
    if isinstance(value, bool):
        return f'<c r="{reference}" t="b"><v>{1 if value else 0}</v></c>'
    if isinstance(value, (int, float)) and not isinstance(value, bool) and math.isfinite(float(value)):
        number = str(value) if isinstance(value, int) else format(float(value), ".15g")
        return f'<c r="{reference}"><v>{number}</v></c>'
    if isinstance(value, (date, datetime)):
        return f'<c r="{reference}" s="1"><v>{format(_excel_date_serial(value), ".15g")}</v></c>'
    text = escape(_text(value))
    preserve = ' xml:space="preserve"' if text[:1].isspace() or text[-1:].isspace() else ""
    return f'<c r="{reference}" t="inlineStr"><is><t{preserve}>{text}</t></is></c>'


def _typed_cell_value(table: dict[str, Any], index: int, value: Any) -> Any:
    """Restore JSON-safe scalar values using the snapshot's declared types."""

    value_types = table.get("value_types")
    kind = value_types[index] if isinstance(value_types, list) and index < len(value_types) else "text"
    if value in (None, ""):
        return value
    if kind == "date" and isinstance(value, str):
        try:
            return datetime.fromisoformat(value.replace("Z", "+00:00")) if "T" in value else date.fromisoformat(value)
        except ValueError:
            return value
    if kind == "number" and isinstance(value, str):
        try:
            number = float(value)
            return int(number) if number.is_integer() else number
        except ValueError:
            return value
    if kind == "boolean" and isinstance(value, str):
        lowered = value.strip().lower()
        if lowered in {"true", "1", "yes", "on"}:
            return True
        if lowered in {"false", "0", "no", "off"}:
            return False
    return value


_XLSX_STYLES_XML = (
    '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
    '<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
    '<numFmts count="1"><numFmt numFmtId="164" formatCode="yyyy-mm-dd hh:mm:ss"/></numFmts>'
    '<fonts count="1"><font><sz val="11"/><name val="Calibri"/></font></fonts>'
    '<fills count="2"><fill><patternFill patternType="none"/></fill>'
    '<fill><patternFill patternType="gray125"/></fill></fills>'
    '<borders count="1"><border><left/><right/><top/><bottom/><diagonal/></border></borders>'
    '<cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs>'
    '<cellXfs count="2">'
    '<xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/>'
    '<xf numFmtId="164" fontId="0" fillId="0" borderId="0" applyNumberFormat="1" xfId="0"/>'
    '</cellXfs>'
    '<cellStyles count="1"><cellStyle name="Normal" xfId="0" builtinId="0"/></cellStyles>'
    '<dxfs count="0"/><tableStyles count="0" defaultTableStyle="TableStyleMedium2" defaultPivotStyle="PivotStyleMedium9"/>'
    '</styleSheet>'
)


def _sheet_xml(rows: list[list[Any]]) -> str:
    rendered_rows: list[str] = []
    for row_index, values in enumerate(rows[:_MAX_ROWS], start=1):
        cells = []
        for column_index, value in enumerate(values[:_MAX_COLUMNS]):
            cells.append(_cell_xml(f"{_column_name(column_index)}{row_index}", value))
        rendered_rows.append(f'<row r="{row_index}">' + "".join(cells) + "</row>")
    return (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        "<sheetData>" + "".join(rendered_rows) + "</sheetData></worksheet>"
    )


def _render_xlsx(envelope: ResultEnvelope, content_shape: str) -> RenderedResult:
    sheets: list[tuple[str, list[list[Any]]]] = []
    if content_shape != "data":
        sheets.append(
            (
                "回答",
                [
                    ["查询", envelope.query],
                    ["回答", envelope.answer],
                    *([["检索概览", envelope.footer]] if envelope.footer else []),
                ],
            )
        )
    for table in envelope.tables if content_shape != "raw" else []:
        columns = list(table.get("columns") or [])
        if not columns:
            continue
        rows = [columns]
        for row in (table.get("rows") or [])[:_MAX_ROWS]:
            values = row if isinstance(row, (list, tuple)) else [row]
            rows.append([
                _typed_cell_value(table, index, values[index]) if index < len(values) else ""
                for index in range(len(columns))
            ])
        sheets.append((str(table.get("name") or "检索结果"), rows))
    if not sheets:
        sheets.append(("结果", [["结果", "未找到可结构化内容"]]))
    if envelope.citations and content_shape != "raw":
        sheets.append(
            (
                "参考来源",
                [
                    ["编号", "证据ID", "来源", "定位", "摘录", "类型"],
                    *[
                        [item.get("index", ""), item.get("evidence_id", ""), item.get("title", ""), item.get("locator", ""), item.get("excerpt", ""), item.get("source_type", "")]
                        for item in envelope.citations[:_MAX_CITATIONS]
                    ],
                ],
            )
        )

    used_names: set[str] = set()
    safe_names = [_safe_sheet_name(name, used_names) for name, _ in sheets]
    workbook_sheets = []
    relationships = []
    content_overrides = []
    for index, name in enumerate(safe_names, start=1):
        workbook_sheets.append(f'<sheet name="{escape(name)}" sheetId="{index}" r:id="rId{index}"/>')
        relationships.append(
            f'<Relationship Id="rId{index}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet{index}.xml"/>'
        )
        content_overrides.append(
            f'<Override PartName="/xl/worksheets/sheet{index}.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>'
        )
    styles_rel_id = len(safe_names) + 1
    relationships.append(
        f'<Relationship Id="rId{styles_rel_id}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>'
    )
    content_overrides.append(
        '<Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>'
    )

    workbook_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships"><sheets>'
        + "".join(workbook_sheets)
        + "</sheets></workbook>"
    )
    workbook_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        + "".join(relationships)
        + '</Relationships>'
    )
    root_rels = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        '<Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>'
        '</Relationships>'
    )
    content_types = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
        '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
        '<Default Extension="xml" ContentType="application/xml"/>'
        '<Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
        + "".join(content_overrides)
        + '</Types>'
    )
    output = io.BytesIO()
    with zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types)
        archive.writestr("_rels/.rels", root_rels)
        archive.writestr("xl/workbook.xml", workbook_xml)
        archive.writestr("xl/_rels/workbook.xml.rels", workbook_rels)
        for index, (_name, rows) in enumerate(sheets, start=1):
            archive.writestr(f"xl/worksheets/sheet{index}.xml", _sheet_xml(rows))
        archive.writestr("xl/styles.xml", _XLSX_STYLES_XML)
    return RenderedResult(
        content=output.getvalue(),
        mime_type=XLSX_MIME,
        extension="xlsx",
        preview={"format": "xlsx", "sheets": [{"name": name, "row_count": len(rows)} for name, rows in zip(safe_names, (rows for _name, rows in sheets))]},
    )


def _table_rows(table: dict[str, Any]) -> tuple[list[str], list[list[Any]]]:
    columns = [str(column) for column in (table.get("columns") or [])][:_MAX_COLUMNS]
    if not columns:
        return [], []
    rows: list[list[Any]] = []
    for row in (table.get("rows") or [])[:_MAX_ROWS]:
        values = row if isinstance(row, (list, tuple)) else [row]
        rows.append([values[index] if index < len(values) else "" for index in range(len(columns))])
    return columns, rows


def _citation_count(envelope: ResultEnvelope) -> int:
    return min(len(envelope.citations), _MAX_CITATIONS)


def _render_docx(envelope: ResultEnvelope, content_shape: str) -> RenderedResult:
    from docx import Document
    from docx.shared import Inches, Pt

    document = Document()
    document.core_properties.title = envelope.title
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.75)
    section.right_margin = Inches(0.75)
    document.add_heading(envelope.title, level=0)
    if content_shape == "report" and envelope.query:
        paragraph = document.add_paragraph()
        paragraph.add_run("查询：").bold = True
        paragraph.add_run(envelope.query)
    if content_shape in {"report", "raw"} and envelope.answer:
        document.add_heading("回答", level=1)
        document.add_paragraph(envelope.answer.rstrip())
    if content_shape == "report" and envelope.footer:
        document.add_heading("检索概览", level=1)
        document.add_paragraph(envelope.footer.rstrip())

    section_count = 0
    for table_data in envelope.tables if content_shape != "raw" else []:
        columns, rows = _table_rows(table_data)
        if not columns:
            continue
        section_count += 1
        document.add_heading(str(table_data.get("name") or "检索结果"), level=1)
        table = document.add_table(rows=1, cols=len(columns))
        table.style = "Table Grid"
        for index, column in enumerate(columns):
            table.rows[0].cells[index].text = column
        for row in rows:
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = _text(value)
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(9)

    if envelope.citations and content_shape != "raw":
        document.add_heading("参考来源", level=1)
        for citation in envelope.citations[:_MAX_CITATIONS]:
            title = _text(citation.get("title") or "参考来源")
            locator = _text(citation.get("locator"))
            suffix = f"（{locator}）" if locator else ""
            paragraph = document.add_paragraph(style="List Number")
            paragraph.add_run(f"{title}{suffix}")
            excerpt = _text(citation.get("excerpt"))
            if excerpt:
                paragraph.add_run(f"\n{excerpt.replace(chr(10), ' ')}")

    output = io.BytesIO()
    document.save(output)
    return RenderedResult(
        content=output.getvalue(),
        mime_type=DOCX_MIME,
        extension="docx",
        preview={
            "format": "docx",
            "section_count": section_count,
            "citation_count": _citation_count(envelope),
        },
    )


def _pdf_paragraph(text: Any, style):
    from reportlab.platypus import Paragraph

    value = escape(_text(text)).replace("\r\n", "\n").replace("\r", "\n").replace("\n", "<br/>")
    return Paragraph(value, style)


def _render_pdf(envelope: ResultEnvelope, content_shape: str) -> RenderedResult:
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_CENTER
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.platypus import SimpleDocTemplate, Spacer, Table, TableStyle

    # STSong-Light is a standard ReportLab CJK font and avoids relying on a
    # server-local browser or a user-provided font path for Chinese results.
    font_name = "STSong-Light"
    try:
        pdfmetrics.registerFont(UnicodeCIDFont(font_name))
    except (KeyError, ValueError):
        # ReportLab may already have registered it in a long-lived worker.
        pass
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "ExportTitle", parent=styles["Title"], fontName=font_name, fontSize=20,
        leading=26, alignment=TA_CENTER, spaceAfter=12,
    )
    heading_style = ParagraphStyle(
        "ExportHeading", parent=styles["Heading2"], fontName=font_name, fontSize=13,
        leading=18, spaceBefore=8, spaceAfter=6,
    )
    body_style = ParagraphStyle(
        "ExportBody", parent=styles["BodyText"], fontName=font_name, fontSize=10.5,
        leading=17, wordWrap="CJK", spaceAfter=6,
    )
    small_style = ParagraphStyle(
        "ExportSmall", parent=body_style, fontSize=8.5, leading=12,
    )
    story = [_pdf_paragraph(envelope.title, title_style)]
    if content_shape == "report" and envelope.query:
        story.extend([_pdf_paragraph(f"查询：{envelope.query}", body_style), Spacer(1, 2 * mm)])
    if content_shape in {"report", "raw"} and envelope.answer:
        story.extend([_pdf_paragraph("回答", heading_style), _pdf_paragraph(envelope.answer.rstrip(), body_style)])
    if content_shape == "report" and envelope.footer:
        story.extend([_pdf_paragraph("检索概览", heading_style), _pdf_paragraph(envelope.footer.rstrip(), body_style)])

    section_count = 0
    for table_data in envelope.tables if content_shape != "raw" else []:
        columns, rows = _table_rows(table_data)
        if not columns:
            continue
        section_count += 1
        story.append(_pdf_paragraph(str(table_data.get("name") or "检索结果"), heading_style))
        table_values = [[_pdf_paragraph(column, small_style) for column in columns]]
        table_values.extend([[_pdf_paragraph(value, small_style) for value in row] for row in rows])
        result_table = Table(table_values, repeatRows=1, hAlign="LEFT")
        result_table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#E8EEF7")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), colors.HexColor("#1F2937")),
                    ("GRID", (0, 0), (-1, -1), 0.35, colors.HexColor("#B7C3D4")),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5),
                    ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                    ("TOPPADDING", (0, 0), (-1, -1), 4),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
                ]
            )
        )
        story.extend([result_table, Spacer(1, 4 * mm)])

    if envelope.citations and content_shape != "raw":
        story.append(_pdf_paragraph("参考来源", heading_style))
        for citation in envelope.citations[:_MAX_CITATIONS]:
            title = _text(citation.get("title") or "参考来源")
            locator = _text(citation.get("locator"))
            suffix = f"（{locator}）" if locator else ""
            line = f"{citation.get('index', '')}. {title}{suffix}"
            excerpt = _text(citation.get("excerpt"))
            if excerpt:
                line += f"\n{excerpt.replace(chr(10), ' ')}"
            story.append(_pdf_paragraph(line, small_style))

    output = io.BytesIO()
    document = SimpleDocTemplate(
        output, pagesize=A4, leftMargin=18 * mm, rightMargin=18 * mm,
        topMargin=16 * mm, bottomMargin=16 * mm, title=envelope.title,
    )
    document.build(story)
    return RenderedResult(
        content=output.getvalue(),
        mime_type=PDF_MIME,
        extension="pdf",
        preview={
            "format": "pdf",
            "page_size": "A4",
            "section_count": section_count,
            "citation_count": _citation_count(envelope),
        },
    )


_PPT_THEMES = {
    "light": {"background": (255, 255, 255), "foreground": (31, 41, 55), "accent": (37, 99, 235)},
    "dark": {"background": (17, 24, 39), "foreground": (243, 244, 246), "accent": (96, 165, 250)},
    "blue": {"background": (239, 246, 255), "foreground": (30, 58, 138), "accent": (29, 78, 216)},
}


def _ppt_color(value: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*value)


def _set_slide_theme(slide, theme: dict[str, tuple[int, int, int]]) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _ppt_color(theme["background"])


def _add_ppt_text(slide, text: Any, *, left: float, top: float, width: float, height: float, font_size: int, color, bold: bool = False):
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = _text(text)
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _ppt_color(color)
    return shape


def _numeric_chart_data(table_data: dict[str, Any]):
    from pptx.chart.data import CategoryChartData

    columns, rows = _table_rows(table_data)
    if len(columns) < 2 or not rows:
        return None
    numeric_index = None
    numeric_values: list[float] = []
    for index in range(1, len(columns)):
        values: list[float] = []
        for row in rows[:10]:
            value = row[index]
            if isinstance(value, bool):
                values = []
                break
            try:
                number = float(value)
            except (TypeError, ValueError):
                values = []
                break
            if number != number or number in {float("inf"), float("-inf")}:
                values = []
                break
            values.append(number)
        if values:
            numeric_index = index
            numeric_values = values
            break
    if numeric_index is None:
        return None
    data = CategoryChartData()
    data.categories = [(_text(row[0]) or str(index + 1)) for index, row in enumerate(rows[:10])]
    data.add_series(columns[numeric_index], numeric_values)
    return data


def _render_pptx(envelope: ResultEnvelope, content_shape: str, render_options: dict[str, Any]) -> RenderedResult:
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData  # noqa: F401 - validates the optional renderer dependency
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches

    theme_name = str(render_options.get("theme") or "light").strip().lower()
    if theme_name not in _PPT_THEMES:
        raise ValueError(f"unsupported presentation theme: {theme_name}")
    theme = _PPT_THEMES[theme_name]
    include_charts = render_options.get("include_charts", True) is True
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    blank = presentation.slide_layouts[6]

    title_slide = presentation.slides.add_slide(blank)
    _set_slide_theme(title_slide, theme)
    _add_ppt_text(title_slide, envelope.title, left=0.7, top=2.1, width=12, height=1.1, font_size=28, color=theme["foreground"], bold=True)
    if content_shape == "report" and envelope.query:
        _add_ppt_text(title_slide, f"查询：{envelope.query}", left=0.75, top=3.5, width=11.8, height=0.8, font_size=16, color=theme["foreground"])

    if content_shape in {"report", "raw"} and (envelope.answer or envelope.footer):
        answer_slide = presentation.slides.add_slide(blank)
        _set_slide_theme(answer_slide, theme)
        _add_ppt_text(answer_slide, "回答", left=0.55, top=0.35, width=12.2, height=0.5, font_size=22, color=theme["accent"], bold=True)
        body = envelope.answer.rstrip()
        if content_shape == "report" and envelope.footer:
            body = f"{body}\n\n检索概览：{envelope.footer.rstrip()}" if body else envelope.footer.rstrip()
        _add_ppt_text(answer_slide, body, left=0.75, top=1.15, width=11.7, height=5.6, font_size=16, color=theme["foreground"])

    chart_count = 0
    section_count = 0
    for table_data in envelope.tables if content_shape != "raw" else []:
        columns, rows = _table_rows(table_data)
        if not columns:
            continue
        section_count += 1
        for chunk_start in range(0, len(rows) or 1, 16):
            chunk = rows[chunk_start:chunk_start + 16]
            slide = presentation.slides.add_slide(blank)
            _set_slide_theme(slide, theme)
            _add_ppt_text(
                slide, str(table_data.get("name") or "检索结果"),
                left=0.55, top=0.3, width=12.2, height=0.5, font_size=20,
                color=theme["accent"], bold=True,
            )
            table = slide.shapes.add_table(len(chunk) + 1, len(columns), Inches(0.55), Inches(1.05), Inches(12.2), Inches(5.7)).table
            for column_index, column in enumerate(columns):
                cell = table.cell(0, column_index)
                cell.text = column
                cell.fill.solid()
                cell.fill.fore_color.rgb = _ppt_color(theme["accent"])
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = _ppt_color((255, 255, 255))
            for row_index, row in enumerate(chunk, start=1):
                for column_index, value in enumerate(row):
                    cell = table.cell(row_index, column_index)
                    cell.text = _text(value)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = _ppt_color(theme["background"])
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = _ppt_color(theme["foreground"])
            if include_charts and chunk_start == 0:
                chart_data = _numeric_chart_data(table_data)
                if chart_data is not None:
                    chart_slide = presentation.slides.add_slide(blank)
                    _set_slide_theme(chart_slide, theme)
                    _add_ppt_text(
                        chart_slide, f"{table_data.get('name') or '检索结果'} - 图表",
                        left=0.55, top=0.3, width=12.2, height=0.5, font_size=20,
                        color=theme["accent"], bold=True,
                    )
                    chart = chart_slide.shapes.add_chart(
                        XL_CHART_TYPE.COLUMN_CLUSTERED,
                        Inches(1.0), Inches(1.15), Inches(11.3), Inches(5.7), chart_data,
                    ).chart
                    chart.has_legend = False
                    chart.chart_title.has_text_frame = True
                    chart.chart_title.text_frame.text = str(table_data.get("name") or "检索结果")
                    chart_count += 1

    if envelope.citations and content_shape != "raw":
        source_slide = presentation.slides.add_slide(blank)
        _set_slide_theme(source_slide, theme)
        _add_ppt_text(source_slide, "参考来源", left=0.55, top=0.3, width=12.2, height=0.5, font_size=20, color=theme["accent"], bold=True)
        source_lines = []
        for citation in envelope.citations[:_MAX_CITATIONS]:
            title = _text(citation.get("title") or "参考来源")
            locator = _text(citation.get("locator"))
            suffix = f"（{locator}）" if locator else ""
            source_lines.append(f"{citation.get('index', '')}. {title}{suffix}")
        _add_ppt_text(source_slide, "\n".join(source_lines), left=0.8, top=1.15, width=11.6, height=5.7, font_size=14, color=theme["foreground"])

    output = io.BytesIO()
    presentation.save(output)
    return RenderedResult(
        content=output.getvalue(),
        mime_type=PPTX_MIME,
        extension="pptx",
        preview={
            "format": "pptx",
            "theme": theme_name,
            "slide_count": len(presentation.slides),
            "section_count": section_count,
            "chart_count": chart_count,
            "citation_count": _citation_count(envelope),
        },
    )


def render_result(
    envelope: ResultEnvelope,
    format: str,
    *,
    content_shape: str = "report",
    render_options: dict[str, Any] | None = None,
) -> RenderedResult:
    shape = normalize_content_shape(content_shape)
    normalized = normalize_export_format(format)
    options = dict(render_options or {})
    if normalized == "md":
        rendered = _render_markdown(envelope, shape)
    elif normalized == "xlsx":
        rendered = _render_xlsx(envelope, shape)
    elif normalized == "docx":
        rendered = _render_docx(envelope, shape)
    elif normalized == "pdf":
        rendered = _render_pdf(envelope, shape)
    elif normalized == "pptx":
        rendered = _render_pptx(envelope, shape, options)
    else:
        raise ValueError(f"unsupported export format: {format}")
    return replace(rendered, preview={**rendered.preview, "content_shape": shape})
